import os
import json
import re
import numpy as np
import faiss
from pypdf import PdfReader
from sentence_transformers import SentenceTransformer
from django.conf import settings

_model = None


def get_embedding_model():
    global _model
    if _model is None:
        _model = SentenceTransformer(settings.RAG_EMBEDDING_MODEL)
    return _model


def extract_text(file_path):
    """Extract text from PDF, DOCX, PPTX, XLSX, or TXT files."""
    lower = file_path.lower()

    if lower.endswith(".pdf"):
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if lower.endswith(".docx"):
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if lower.endswith(".pptx"):
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
        return "\n".join(texts)

    if lower.endswith(".xlsx"):
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append("\t".join(cells))
        wb.close()
        return "\n".join(rows)

    # Fallback: plain text
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def chunk_text(text, chunk_size=None, overlap=None):
    """Split text into overlapping chunks at sentence boundaries."""
    chunk_size = chunk_size or settings.RAG_CHUNK_SIZE
    overlap = overlap or settings.RAG_CHUNK_OVERLAP

    # Normalize PDF line breaks: single newlines (column wrapping) → spaces,
    # but preserve double+ newlines as paragraph separators
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
    # Collapse runs of whitespace (except paragraph breaks)
    text = re.sub(r'[^\S\n]+', ' ', text)

    # Split into sentences (period/question/exclamation followed by space, or paragraph break)
    sentences = re.split(r'(?<=[.!?])\s+|\n\n+', text)
    sentences = [s.strip() for s in sentences if s.strip()]

    chunks = []
    current_chunk = ""

    for sentence in sentences:
        if current_chunk and len(current_chunk) + len(sentence) + 1 > chunk_size:
            chunks.append(current_chunk)
            words = current_chunk.split()
            overlap_text = " ".join(words[-overlap // 5:]) if len(words) > overlap // 5 else current_chunk
            current_chunk = overlap_text + " " + sentence
        else:
            current_chunk = (current_chunk + " " + sentence).strip() if current_chunk else sentence

    if current_chunk:
        chunks.append(current_chunk)

    return chunks


def _user_index_dir(user_id):
    """Return the directory for a user's FAISS index."""
    path = os.path.join(settings.RAG_INDEX_DIR, str(user_id))
    os.makedirs(path, exist_ok=True)
    return path


def _load_index(user_id):
    """Load a user's FAISS index and metadata. Creates new if not found."""
    index_dir = _user_index_dir(user_id)
    index_path = os.path.join(index_dir, "index.faiss")
    meta_path = os.path.join(index_dir, "metadata.json")

    if os.path.exists(index_path) and os.path.exists(meta_path):
        index = faiss.read_index(index_path)
        with open(meta_path, "r") as f:
            metadata = json.load(f)
    else:
        dim = get_embedding_model().get_sentence_embedding_dimension()
        index = faiss.IndexFlatL2(dim)
        metadata = []

    return index, metadata


def _save_index(user_id, index, metadata):
    """Persist a user's FAISS index and metadata to disk."""
    index_dir = _user_index_dir(user_id)
    faiss.write_index(index, os.path.join(index_dir, "index.faiss"))
    with open(os.path.join(index_dir, "metadata.json"), "w") as f:
        json.dump(metadata, f)


def index_document(user_id, filename, text):
    """Chunk, embed, and add a document to the user's FAISS index."""
    chunks = chunk_text(text)
    if not chunks:
        return 0

    model = get_embedding_model()
    embeddings = model.encode(chunks, normalize_embeddings=True)
    embeddings = np.array(embeddings, dtype="float32")

    index, metadata = _load_index(user_id)

    # Remove old chunks for this document so re-uploads replace, not duplicate
    keep = [i for i, m in enumerate(metadata) if m["filename"] != filename]
    if len(keep) < len(metadata):
        if keep:
            kept_vecs = np.array(
                [index.reconstruct(i) for i in keep], dtype="float32"
            )
            new_index = faiss.IndexFlatL2(kept_vecs.shape[1])
            new_index.add(kept_vecs)
        else:
            new_index = faiss.IndexFlatL2(embeddings.shape[1])
        metadata = [metadata[i] for i in keep]
        index = new_index

    index.add(embeddings)

    for chunk in chunks:
        metadata.append({"filename": filename, "text": chunk})

    _save_index(user_id, index, metadata)
    return len(chunks)


def retrieve(user_id, query, top_k=None):
    """Hybrid retrieval: semantic similarity + keyword re-ranking."""
    top_k = top_k or settings.RAG_TOP_K
    index, metadata = _load_index(user_id)

    if index.ntotal == 0:
        return []

    model = get_embedding_model()
    query_vec = model.encode([query], normalize_embeddings=True)
    query_vec = np.array(query_vec, dtype="float32")

    # Fetch extra candidates for re-ranking
    n_candidates = min(top_k * 3, index.ntotal)
    distances, indices = index.search(query_vec, n_candidates)

    # Extract meaningful query terms for keyword boosting
    stop_words = {
        "the", "a", "an", "is", "are", "was", "were", "what", "which",
        "who", "how", "does", "did", "and", "or", "but", "in", "on",
        "at", "to", "for", "of", "with", "by", "from", "this", "that",
        "it", "its", "about", "talk", "talks", "can", "you", "me",
    }
    query_terms = [
        w for w in re.split(r'\W+', query.lower())
        if len(w) > 2 and w not in stop_words
    ]

    scored = []
    for rank, idx in enumerate(indices[0]):
        if idx < 0 or idx >= len(metadata):
            continue
        chunk_lower = metadata[idx]["text"].lower()
        keyword_hits = sum(1 for t in query_terms if t in chunk_lower)
        semantic_score = 1.0 / (1.0 + float(distances[0][rank]))
        keyword_score = keyword_hits / max(len(query_terms), 1)
        combined = semantic_score + keyword_score * 0.5
        scored.append((combined, idx))

    scored.sort(key=lambda x: x[0], reverse=True)

    results = []
    for _, idx in scored[:top_k]:
        results.append(metadata[idx])
    return results

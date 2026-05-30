import os
import re

import numpy as np
import faiss
from openai import OpenAI
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from django.conf import settings

from .embeddings import embed_texts
from .vector_store import load_index, save_index


def extract_text(file_path):
    """Extract text from PDF, DOCX, PPTX, XLSX, or TXT files."""
    lower = file_path.lower()

    if lower.endswith(".pdf"):
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if lower.endswith(".docx"):
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if lower.endswith(".pptx"):
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
        return "\n".join(texts)

    if lower.endswith(".xlsx"):
        wb = load_workbook(file_path, read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    rows.append("\t".join(cells))
        wb.close()
        return "\n".join(rows)

    # Fallback: plain text
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def chunk_text(text, chunk_size=None, overlap=None):
    """Split text into overlapping chunks at sentence boundaries."""
    chunk_size = chunk_size or settings.RAG_CHUNK_SIZE
    overlap = overlap or settings.RAG_CHUNK_OVERLAP

    # Normalize PDF line breaks: single newlines (column wrapping) -> spaces,
    # but preserve double+ newlines as paragraph separators
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)
    # Collapse runs of whitespace (except paragraph breaks)
    text = re.sub(r'[^\S\n]+', ' ', text)

    # Split into sentences (period/question/exclamation followed by space, or paragraph break)
    sentences = re.split(r'(?<=[.!?])\s+|\n\n+', text)
    sentences = [s.strip() for s in sentences if s.strip()]

    chunks = []
    current_chunk = ""

    for sentence in sentences:
        if current_chunk and len(current_chunk) + len(sentence) + 1 > chunk_size:
            chunks.append(current_chunk)
            words = current_chunk.split()
            overlap_text = " ".join(words[-overlap // 5:]) if len(words) > overlap // 5 else current_chunk
            current_chunk = overlap_text + " " + sentence
        else:
            current_chunk = (current_chunk + " " + sentence).strip() if current_chunk else sentence

    if current_chunk:
        chunks.append(current_chunk)

    return chunks


def index_document(user_id, filename, text):
    """Chunk, embed, and add a document to the user's FAISS index."""
    chunks = chunk_text(text)
    if not chunks:
        return 0

    embeddings = embed_texts(chunks)

    index, metadata = load_index(user_id)

    # Remove old chunks for this document so re-uploads replace, not duplicate
    keep = [i for i, m in enumerate(metadata) if m["filename"] != filename]
    if len(keep) < len(metadata):
        if keep:
            kept_vecs = np.array(
                [index.reconstruct(i) for i in keep], dtype="float32"
            )
            new_index = faiss.IndexFlatL2(kept_vecs.shape[1])
            new_index.add(kept_vecs)
        else:
            new_index = faiss.IndexFlatL2(embeddings.shape[1])
        metadata = [metadata[i] for i in keep]
        index = new_index

    index.add(embeddings)

    for chunk in chunks:
        metadata.append({"filename": filename, "text": chunk})

    save_index(user_id, index, metadata)
    return len(chunks)


_chat_client = None


def _get_chat_client():
    global _chat_client
    if _chat_client is None:
        _chat_client = OpenAI(api_key=settings.OPENAI_API_KEY)
    return _chat_client


def _rewrite_query(query):
    """Repair typos and light-paraphrase a user query before retrieval.

    The original query still flows through to the answer-generating LLM call;
    only retrieval (embedding, keyword scoring, filename matching) uses the
    rewritten version. Returns the original query on any failure so retrieval
    never blocks on the rewrite call.
    """
    try:
        response = _get_chat_client().chat.completions.create(
            model=settings.OPENAI_MODEL,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You preprocess search queries for a document "
                        "retrieval system. Rewrite the user's question to: "
                        "(1) fix obvious spelling typos, "
                        "(2) preserve filenames, proper nouns, and technical "
                        "terms exactly, "
                        "(3) keep the intent and length similar. "
                        "Return only the rewritten query, no explanation, "
                        "no quotes."
                    ),
                },
                {"role": "user", "content": query},
            ],
            max_tokens=200,
            temperature=0,
        )
        rewritten = response.choices[0].message.content.strip().strip('"').strip("'")
        return rewritten or query
    except Exception:
        return query


def retrieve(user_id, query, top_k=None):
    """Hybrid retrieval: semantic similarity + keyword re-ranking +
    filename-mention boost. Queries pass through an LLM typo-repair step
    before being embedded.
    """
    top_k = top_k or settings.RAG_TOP_K
    index, metadata = load_index(user_id)

    if index.ntotal == 0:
        return []

    # Repair typos before retrieval. Shadowing `query` is intentional —
    # everything below (embedding, filename detection, keyword scoring)
    # should use the cleaned version.
    query = _rewrite_query(query)

    query_vec = embed_texts([query])

    # Score every chunk in the index, not just the top-N nearest. The diversity
    # pass below can only surface a doc if at least one of its chunks reaches
    # this scoring loop — capping candidates at top_k*3 lets a chunk-heavy doc
    # drown out a chunk-light one. Per-user indexes are small (≪10k chunks)
    # so an exhaustive scan is cheap.
    n_candidates = index.ntotal
    distances, indices = index.search(query_vec, n_candidates)

    # Extract meaningful query terms for keyword boosting
    stop_words = {
        "the", "a", "an", "is", "are", "was", "were", "what", "which",
        "who", "how", "does", "did", "and", "or", "but", "in", "on",
        "at", "to", "for", "of", "with", "by", "from", "this", "that",
        "it", "its", "about", "talk", "talks", "can", "you", "me",
    }
    query_terms = [
        w for w in re.split(r'\W+', query.lower())
        if len(w) > 2 and w not in stop_words
    ]

    # If the query mentions a word stem from any indexed filename, treat that
    # file as the user's target and boost its chunks above the field. Without
    # this, a question like "what's in foo.pdf" can't outrank a larger doc
    # that semantically matches "what's in" generically.
    query_lower = query.lower()
    ext_stems = {"pdf", "docx", "pptx", "xlsx", "txt", "doc", "ppt", "xls"}
    targeted_files = set()
    for m in metadata:
        base = os.path.basename(m["filename"]).lower()
        stems = [s for s in re.findall(r'[a-z]{4,}', base) if s not in ext_stems]
        if any(s in query_lower for s in stems):
            targeted_files.add(m["filename"])

    # When the user names a file, they want answers from that file regardless
    # of which section the answer lives in or how badly the question is typo'd.
    # Expand top_k so every chunk of the targeted file fits, plus one per other
    # doc for cross-doc context. Capped so huge docs don't blow the context.
    MAX_TOP_K = 50
    if targeted_files:
        targeted_chunk_count = sum(
            1 for m in metadata if m["filename"] in targeted_files
        )
        other_doc_count = len(
            {m["filename"] for m in metadata} - targeted_files
        )
        top_k = min(max(top_k, targeted_chunk_count + other_doc_count), MAX_TOP_K)

    scored = []
    for rank, idx in enumerate(indices[0]):
        if idx < 0 or idx >= len(metadata):
            continue
        chunk_lower = metadata[idx]["text"].lower()
        keyword_hits = sum(1 for t in query_terms if t in chunk_lower)
        semantic_score = 1.0 / (1.0 + float(distances[0][rank]))
        keyword_score = keyword_hits / max(len(query_terms), 1)
        # filename_boost dominates the other terms (max semantic≈1, max
        # keyword*0.5=0.5) so a targeted file's chunks always rank above
        # non-targeted ones. Combined with the top_k bump above, this means
        # every targeted chunk gets pulled in before any non-targeted one.
        filename_boost = 2.0 if metadata[idx]["filename"] in targeted_files else 0.0
        combined = semantic_score + keyword_score * 0.5 + filename_boost
        scored.append((combined, idx))

    scored.sort(key=lambda x: x[0], reverse=True)

    # Drop chunks far below the top score before the diversity pass. Without
    # this floor, the per-document pass below pulled one chunk from EVERY
    # indexed doc — even off-topic ones — and their filenames then showed
    # up in `sources`. 0.6 keeps near-equal matches and excludes clear
    # misses; the top chunk itself always survives (top * 0.6 < top).
    top_score = scored[0][0] if scored else 0.0
    relevance_floor = top_score * 0.6

    # Diversity: one chunk per relevant document first, then fill the rest.
    doc_best = {}
    for score, idx in scored:
        doc = metadata[idx]["filename"]
        if doc not in doc_best:
            doc_best[doc] = (score, idx)

    selected = set()
    results = []

    for score, idx in doc_best.values():
        if score < relevance_floor:
            continue
        if len(results) >= top_k:
            break
        results.append((score, idx))
        selected.add(idx)

    # scored is sorted desc, so the first sub-floor chunk means we're done.
    for score, idx in scored:
        if score < relevance_floor:
            break
        if len(results) >= top_k:
            break
        if idx not in selected:
            results.append((score, idx))
            selected.add(idx)

    results.sort(key=lambda x: x[0], reverse=True)
    return [metadata[idx] for _, idx in results]


def build_rag_conversation(content, results):
    """Build the RAG-mode conversation from retrieved chunks.

    Returns (conversation, sources, retrieved_chunks).
    """
    sources = list(dict.fromkeys(r["filename"] for r in results))
    retrieved_chunks = [r["text"] for r in results]

    # Label each chunk with its source document
    labeled = []
    for r in results:
        doc_name = os.path.basename(r["filename"])
        labeled.append(f"[Document: {doc_name}]\n{r['text']}")
    context_text = "\n\n---\n\n".join(labeled)

    # RAG: answer-focused prompt. Refusal threshold is intentionally loose —
    # we'd rather have the model attempt an answer from related context than
    # bail out because the user typed "absrtact" instead of "abstract".
    conversation = [
        {
            "role": "system",
            "content": (
                "You are a helpful research assistant. Answer the user's "
                "question using only the document excerpts provided below. "
                "The user's question may contain typos or paraphrase the "
                "document — match by intent, not exact spelling. "
                "Include direct quotes in quotation marks, name the document "
                "each quote came from, then explain in your own words. "
                "Do not use outside knowledge. "
                "If the excerpts cover the topic but don't contain the "
                "specific section or detail the user asked for, answer with "
                "what the excerpts do say and note what's missing. "
                "Respond with \"I don't know. The answer to your question "
                "was not found in the uploaded documents.\" ONLY when none "
                "of the excerpts plausibly relate to the question's topic."
            ),
        },
        {
            "role": "user",
            "content": (
                f"Here are excerpts from my uploaded documents:\n\n"
                f"{context_text}\n\n"
                f"Based on ALL the above excerpts, answer this question:\n"
                f"{content}\n\n"
                f"Include relevant direct quotes from each document "
                f"(mention the document name), then explain what they mean."
            ),
        },
    ]
    return conversation, sources, retrieved_chunks
